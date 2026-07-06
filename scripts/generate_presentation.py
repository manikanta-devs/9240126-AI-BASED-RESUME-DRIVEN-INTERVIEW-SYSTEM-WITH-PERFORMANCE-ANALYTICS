import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set presentation size to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # Color Palette Constants
    # -------------------------------------------------------------
    DARK_BG = RGBColor(15, 23, 42)       # Slate 900 (#0F172A)
    LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50 (#F8FAFC)
    
    TEXT_DARK_BG = RGBColor(248, 250, 252)  # Soft White
    TEXT_MUTED_DARK = RGBColor(148, 163, 184) # Muted Slate
    TEXT_LIGHT_BG = RGBColor(15, 23, 42)   # Slate 900
    TEXT_MUTED_LIGHT = RGBColor(71, 85, 105) # Slate 600
    
    ACCENT_TEAL = RGBColor(13, 148, 136)   # Teal 600 (#0D9488)
    ACCENT_CYAN = RGBColor(6, 182, 212)    # Cyan 500 (#06B6D4)
    ACCENT_BLUE = RGBColor(37, 99, 235)    # Blue 600 (#2563EB)
    
    CARD_BG_LIGHT = RGBColor(255, 255, 255) # Pure White
    CARD_BORDER_LIGHT = RGBColor(226, 232, 240) # Slate 200 (#E2E8F0)
    
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # Helper Functions
    # -------------------------------------------------------------
    def apply_bg(slide, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        return rect

    def add_header(slide, title_text, is_dark_bg=False):
        # Title text box
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Georgia"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK_BG if is_dark_bg else TEXT_LIGHT_BG
        
        # Divider Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.3), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_CYAN if is_dark_bg else ACCENT_TEAL
        line.line.fill.background()

    def draw_card(slide, left, top, width, height, bg_color, border_color=None, border_width=1.0):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(border_width)
        else:
            card.line.fill.background()
        return card

    def add_paragraph(tf, text, font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=False, space_after=6, bullet=False):
        if tf.paragraphs and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(space_after)
        if bullet:
            p.level = 0
        return p

    # -------------------------------------------------------------
    # SLIDE 1: Title (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_bg(slide1, DARK_BG)
    
    # Left accent colored strip
    strip = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT_CYAN
    strip.line.fill.background()
    
    # Large Title + Subtitle
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p_title = add_paragraph(tf1, "AI-Based Resume Driven\nInterview System", font_name="Georgia", font_size=46, color=TEXT_DARK_BG, bold=True, space_after=12)
    p_sub = add_paragraph(tf1, "Intelligent Full-Stack Practice Platform with Performance Analytics", font_name="Calibri", font_size=20, color=ACCENT_CYAN, space_after=24)
    
    # Divider line on Title slide
    line1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.3), Inches(3.0), Inches(0.04))
    line1.fill.solid()
    line1.fill.fore_color.rgb = ACCENT_TEAL
    line1.line.fill.background()
    
    # Prepared by details
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(8.0), Inches(1.5))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    add_paragraph(tf_info, "Prepared by: Project Team", font_name="Calibri", font_size=15, color=TEXT_DARK_BG, bold=True, space_after=4)
    add_paragraph(tf_info, "Powered by React, Flask, Google Gemini AI & spaCy NLP", font_name="Calibri", font_size=13, color=TEXT_MUTED_DARK)

    # -------------------------------------------------------------
    # SLIDE 2: Project Objectives (Light Theme)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide2, LIGHT_BG)
    add_header(slide2, "Project Objective & Mission")
    
    # Main objective card
    draw_card(slide2, Inches(0.75), Inches(1.8), Inches(11.833), Inches(1.2), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    obj_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.333), Inches(1.0))
    tf_obj = obj_box.text_frame
    tf_obj.word_wrap = True
    add_paragraph(tf_obj, "Mission Statement:", font_name="Georgia", font_size=16, color=ACCENT_TEAL, bold=True, space_after=4)
    add_paragraph(tf_obj, "To build an intelligent, responsive platform that bridges the gap between static resume preparation and real-world interviews. By generating tailored questions and assessing text and voice responses in real time, the platform delivers personalized, measurable coaching feedback.", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG)
    
    # Three Pillar Cards side-by-side
    pillars = [
        {"title": "🎯 Personalized Questions", "desc": "Analyzes the user's resume using spaCy NLP and generates role-specific technical and behavioral questions using Gemini AI."},
        {"title": "🗣️ Communication Coach", "desc": "Evaluates candidate responses for speaking pace, clarity, confidence, and flags filler-word patterns to build real interview readiness."},
        {"title": "📊 Performance Analytics", "desc": "Aggregates interview session scores over time. Displays weakness analyses, skill radar charts, and comparison statistics."}
    ]
    
    card_width = Inches(3.7)
    card_gap = Inches(0.366)
    left_margin = Inches(0.75)
    
    for i, pillar in enumerate(pillars):
        card_left = left_margin + i * (card_width + card_gap)
        draw_card(slide2, card_left, Inches(3.3), card_width, Inches(3.2), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
        
        tb = slide2.shapes.add_textbox(card_left + Inches(0.25), Inches(3.55), card_width - Inches(0.5), Inches(2.7))
        tf = tb.text_frame
        tf.word_wrap = True
        add_paragraph(tf, pillar["title"], font_name="Georgia", font_size=18, color=TEXT_LIGHT_BG, bold=True, space_after=12)
        add_paragraph(tf, pillar["desc"], font_name="Calibri", font_size=14, color=TEXT_MUTED_LIGHT)

    # -------------------------------------------------------------
    # SLIDE 3: Problem Statement (Light Theme)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_bg(slide3, LIGHT_BG)
    add_header(slide3, "Problem Statement")
    
    # Left Box: The Current Situation
    draw_card(slide3, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.7), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb_prob_left = slide3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.3))
    tf_pl = tb_prob_left.text_frame
    tf_pl.word_wrap = True
    add_paragraph(tf_pl, "The Challenges Candidates Face", font_name="Georgia", font_size=18, color=ACCENT_TEAL, bold=True, space_after=16)
    
    add_paragraph(tf_pl, "• Generic Question Banks:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_pl, "  Most prep tools provide static questions that do not align with a candidate's projects or skills.", font_name="Calibri", font_size=14, color=TEXT_MUTED_LIGHT, space_after=12)
    
    add_paragraph(tf_pl, "• Absence of Constructive Feedback:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_pl, "  Candidates cannot identify specific reasons why their answers lack depth, technical correctness, or clear delivery.", font_name="Calibri", font_size=14, color=TEXT_MUTED_LIGHT, space_after=12)
    
    add_paragraph(tf_pl, "• No Structured Progress Metrics:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_pl, "  Lack of interactive dashboards makes it difficult for users to track progress and identify persistent weak topics.", font_name="Calibri", font_size=14, color=TEXT_MUTED_LIGHT)

    # Right Box: Why It Matters (High-contrast Accent Card)
    draw_card(slide3, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.7), DARK_BG)
    tb_prob_right = slide3.shapes.add_textbox(Inches(7.28), Inches(2.4), Inches(5.0), Inches(3.5))
    tf_pr = tb_prob_right.text_frame
    tf_pr.word_wrap = True
    add_paragraph(tf_pr, "The Opportunity", font_name="Georgia", font_size=24, color=ACCENT_CYAN, bold=True, space_after=16)
    add_paragraph(tf_pr, "Traditional interview prep relies heavily on manual mock interviews, which are expensive, logistically complex, and subjective.", font_name="Calibri", font_size=16, color=TEXT_DARK_BG, space_after=14)
    add_paragraph(tf_pr, "By leveraging Large Language Models (LLMs) and local NLP, we can democratize premium, personalized, and objective interview practice at zero marginal cost.", font_name="Calibri", font_size=16, color=TEXT_MUTED_DARK)

    # -------------------------------------------------------------
    # SLIDE 4: Literature Survey & Core Research (Light Theme)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_bg(slide4, LIGHT_BG)
    add_header(slide4, "Literature Survey & Technological Foundation")
    
    # 2x2 grid of research domains
    domains = [
        {
            "title": "1. Resume Parsing & NER",
            "points": [
                "Transitioned from rule-based parsers to Named Entity Recognition (NER).",
                "Uses spaCy packages to extract skills, project domains, and roles.",
                "Determines structural entities for dynamic prompt engineering."
            ]
        },
        {
            "title": "2. Adaptive Question Gen",
            "points": [
                "Replaces hardcoded question trees with Prompt Engineering.",
                "Employs system-instruction prompts with structured JSON schemas.",
                "Integrates offline fallback datasets for resilience."
            ]
        },
        {
            "title": "3. Automatic Speech & Evaluation",
            "points": [
                "Utilizes Web Speech API for low-latency voice capture.",
                "Evaluates answers on 4 key metrics: accuracy, depth, relevance, and clarity.",
                "Dual-scoring: LLM API backed by heuristic validation rules."
            ]
        },
        {
            "title": "4. Learning Analytics Dashboards",
            "points": [
                "Relies on visual feedback loops to maximize learning.",
                "Leverages radar charts, performance trends, and weak area tables.",
                "Keeps records persistently using structured JSON documents."
            ]
        }
    ]
    
    grid_w = Inches(5.6)
    grid_h = Inches(2.1)
    
    # Box 1
    draw_card(slide4, Inches(0.75), Inches(1.8), grid_w, grid_h, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb = slide4.shapes.add_textbox(Inches(0.9), Inches(1.9), grid_w - Inches(0.3), grid_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_paragraph(tf, domains[0]["title"], font_name="Georgia", font_size=16, color=ACCENT_TEAL, bold=True, space_after=6)
    for pt in domains[0]["points"]:
        add_paragraph(tf, "• " + pt, font_name="Calibri", font_size=13, color=TEXT_LIGHT_BG, space_after=3)
        
    # Box 2
    draw_card(slide4, Inches(6.98), Inches(1.8), grid_w, grid_h, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb = slide4.shapes.add_textbox(Inches(7.13), Inches(1.9), grid_w - Inches(0.3), grid_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_paragraph(tf, domains[1]["title"], font_name="Georgia", font_size=16, color=ACCENT_TEAL, bold=True, space_after=6)
    for pt in domains[1]["points"]:
        add_paragraph(tf, "• " + pt, font_name="Calibri", font_size=13, color=TEXT_LIGHT_BG, space_after=3)

    # Box 3
    draw_card(slide4, Inches(0.75), Inches(4.2), grid_w, grid_h, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb = slide4.shapes.add_textbox(Inches(0.9), Inches(4.3), grid_w - Inches(0.3), grid_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_paragraph(tf, domains[2]["title"], font_name="Georgia", font_size=16, color=ACCENT_TEAL, bold=True, space_after=6)
    for pt in domains[2]["points"]:
        add_paragraph(tf, "• " + pt, font_name="Calibri", font_size=13, color=TEXT_LIGHT_BG, space_after=3)

    # Box 4
    draw_card(slide4, Inches(6.98), Inches(4.2), grid_w, grid_h, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb = slide4.shapes.add_textbox(Inches(7.13), Inches(4.3), grid_w - Inches(0.3), grid_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_paragraph(tf, domains[3]["title"], font_name="Georgia", font_size=16, color=ACCENT_TEAL, bold=True, space_after=6)
    for pt in domains[3]["points"]:
        add_paragraph(tf, "• " + pt, font_name="Calibri", font_size=13, color=TEXT_LIGHT_BG, space_after=3)

    # -------------------------------------------------------------
    # SLIDE 5: Key Features (Light Theme)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_bg(slide5, LIGHT_BG)
    add_header(slide5, "Key System Features")
    
    # 3x2 grid of small cards
    features = [
        {"title": "📄 Resume Parser & Grade", "desc": "Extracts skills, experience, education; grades completeness score."},
        {"title": "🤖 AI-Driven Question Gen", "desc": "Generates 3-10 questions dynamically tailored to resume + target role."},
        {"title": "🎙️ Multi-Modal Interview Flow", "desc": "Live countdown timer, support for text, speech, and video capture."},
        {"title": "⚡ Answer Evaluator", "desc": "Generates detailed breakdown (Technical Accuracy, Relevance, Clarity, Depth)."},
        {"title": "🗣️ Communication Coach", "desc": "Speaking drills, rate monitoring, filler word analysis, and structure coaching."},
        {"title": "📊 Interactive Dashboards", "desc": "Track session history, radar charts, weakness areas, and growth trends."}
    ]
    
    fw = Inches(3.7)
    fh = Inches(2.1)
    
    col_coords = [Inches(0.75), Inches(4.816), Inches(8.883)]
    row_coords = [Inches(1.8), Inches(4.2)]
    
    for idx, feat in enumerate(features):
        col = idx % 3
        row = idx // 3
        
        c_left = col_coords[col]
        c_top = row_coords[row]
        
        draw_card(slide5, c_left, c_top, fw, fh, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
        
        tb = slide5.shapes.add_textbox(c_left + Inches(0.2), c_top + Inches(0.15), fw - Inches(0.4), fh - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        add_paragraph(tf, feat["title"], font_name="Georgia", font_size=16, color=ACCENT_TEAL, bold=True, space_after=6)
        add_paragraph(tf, feat["desc"], font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT)

    # -------------------------------------------------------------
    # SLIDE 6: Tech Stack & Architecture (Light Theme)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    apply_bg(slide6, LIGHT_BG)
    add_header(slide6, "System Architecture & Tech Stack")
    
    stack = [
        {
            "layer": "Frontend Layer",
            "color": ACCENT_BLUE,
            "techs": [
                "React 18 + Vite - High speed SPA framework",
                "Tailwind CSS - Responsive utilities & modern UI",
                "Recharts - Radar, line, and bar chart visuals",
                "Axios - Handles backend API communications",
                "React Context API - Manages global state"
            ]
        },
        {
            "layer": "Backend Layer",
            "color": ACCENT_TEAL,
            "techs": [
                "Python 3.10+ & Flask 3 - Core application server",
                "spaCy (en_core_web_sm) - Local NLP parsing",
                "PyPDF2 & python-docx - File text extraction",
                "Pydantic - Robust schemas & request validation",
                "Flask-Limiter - Rate-limiting API security"
            ]
        },
        {
            "layer": "AI & Provider Chain",
            "color": ACCENT_CYAN,
            "techs": [
                "Google Gemini 1.5 Flash - Core LLM engine",
                "Hugging Face Inference API - Primary QA model",
                "Fallback Heuristic Engine - Scores local results offline",
                "JSON Parsing Wrappers - Ensures schema compliance",
                "Dual Mode - Auto-switches to local when offline"
            ]
        }
    ]
    
    sw = Inches(3.7)
    sh = Inches(4.7)
    
    for idx, s in enumerate(stack):
        c_left = col_coords[idx]
        
        # Header block for the card
        draw_card(slide6, c_left, Inches(1.8), sw, sh, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
        
        # Mini title banner on top of the card
        banner = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, Inches(1.8), sw, Inches(0.6))
        banner.fill.solid()
        banner.fill.fore_color.rgb = s["color"]
        banner.line.fill.background()
        
        tb_ban = slide6.shapes.add_textbox(c_left, Inches(1.9), sw, Inches(0.4))
        tf_ban = tb_ban.text_frame
        p_ban = tf_ban.paragraphs[0]
        p_ban.text = s["layer"]
        p_ban.alignment = PP_ALIGN.CENTER
        p_ban.font.name = "Georgia"
        p_ban.font.size = Pt(16)
        p_ban.font.bold = True
        p_ban.font.color.rgb = TEXT_DARK_BG
        
        tb = slide6.shapes.add_textbox(c_left + Inches(0.2), Inches(2.6), sw - Inches(0.4), sh - Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        for tech in s["techs"]:
            parts = tech.split(" - ")
            p = add_paragraph(tf, "• " + parts[0], font_name="Calibri", font_size=14, color=TEXT_LIGHT_BG, bold=True, space_after=1, bullet=False)
            add_paragraph(tf, "  " + parts[1], font_name="Calibri", font_size=12, color=TEXT_MUTED_LIGHT, space_after=8)

    # -------------------------------------------------------------
    # SLIDE 7: Methodology & System Flow (Light Theme)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    apply_bg(slide7, LIGHT_BG)
    add_header(slide7, "System Methodology & User Flow")
    
    flow_steps = [
        {"num": "1", "title": "Resume Parse", "desc": "User uploads resume. spaCy parses skills, education, and role data."},
        {"num": "2", "title": "Dynamic Prompt", "desc": "Combines target job description and parsed skills into a custom system prompt."},
        {"num": "3", "title": "Interactive Interview", "desc": "User submits text or voice responses. System records and processes inputs."},
        {"num": "4", "title": "AI Assessment", "desc": "Evaluates response correctness, filler count, and completeness."},
        {"num": "5", "title": "Visual Dashboard", "desc": "Updates charts, tracks progress over time, and highlights weak topics."}
    ]
    
    step_w = Inches(2.1)
    step_h = Inches(4.0)
    gap = Inches(0.33)
    start_left = Inches(0.75)
    
    for idx, step in enumerate(flow_steps):
        c_left = start_left + idx * (step_w + gap)
        
        # Base Card
        draw_card(slide7, c_left, Inches(2.2), step_w, step_h, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
        
        # Step Number Badge
        badge = slide7.shapes.add_shape(MSO_SHAPE.OVAL, c_left + Inches(0.75), Inches(1.8), Inches(0.6), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_TEAL
        badge.line.fill.background()
        
        tb_b = slide7.shapes.add_textbox(c_left + Inches(0.75), Inches(1.85), Inches(0.6), Inches(0.6))
        p_b = tb_b.text_frame.paragraphs[0]
        p_b.text = step["num"]
        p_b.alignment = PP_ALIGN.CENTER
        p_b.font.name = "Georgia"
        p_b.font.size = Pt(16)
        p_b.font.bold = True
        p_b.font.color.rgb = TEXT_DARK_BG
        
        # Text Content inside Card
        tb = slide7.shapes.add_textbox(c_left + Inches(0.15), Inches(2.6), step_w - Inches(0.3), step_h - Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        add_paragraph(tf, step["title"], font_name="Georgia", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=8)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_paragraph(tf, step["desc"], font_name="Calibri", font_size=12, color=TEXT_MUTED_LIGHT)
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 8: Implementation Details (Light Theme)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    apply_bg(slide8, LIGHT_BG)
    add_header(slide8, "Work Completed & Technical Highlights")
    
    # Column 1: Backend & Core services
    draw_card(slide8, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.7), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb_imp_left = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.3))
    tf_imp_l = tb_imp_left.text_frame
    tf_imp_l.word_wrap = True
    add_paragraph(tf_imp_l, "Backend Modules Implemented", font_name="Georgia", font_size=18, color=ACCENT_TEAL, bold=True, space_after=12)
    
    add_paragraph(tf_imp_l, "• Dual-Mode AI Orchestrator:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=1)
    add_paragraph(tf_imp_l, "  `gemini_service.py` manages API retries and falls back seamlessly to Hugging Face or deterministic scripts.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=8)
    
    add_paragraph(tf_imp_l, "• Custom NLP Resume Parsing:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=1)
    add_paragraph(tf_imp_l, "  `resume_service.py` processes PDF/DOCX and isolates skills, experience levels, and roles using spaCy NER matcher.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=8)
    
    add_paragraph(tf_imp_l, "• Comprehensive Unit Testing:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=1)
    add_paragraph(tf_imp_l, "  Configured `pytest` for configs, exception states, validation modules, and endpoints with full HTML cov reports.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT)

    # Column 2: Frontend & Deployment
    draw_card(slide8, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.7), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb_imp_right = slide8.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.3))
    tf_imp_r = tb_imp_right.text_frame
    tf_imp_r.word_wrap = True
    add_paragraph(tf_imp_r, "Frontend & Security Deployments", font_name="Georgia", font_size=18, color=ACCENT_TEAL, bold=True, space_after=12)
    
    add_paragraph(tf_imp_r, "• Responsive SPA Dashboard:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=1)
    add_paragraph(tf_imp_r, "  Integrated React Router, AppContext, Sidebar navigation, Dark Mode support, and real-time dashboard analytics.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=8)
    
    add_paragraph(tf_imp_r, "• API Rate Limiting & Security Headers:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=1)
    add_paragraph(tf_imp_r, "  Uses Flask-Limiter to throttle abusers. Sets strict CSP headers, XSS protections, and limits upload size to 16MB.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=8)
    
    add_paragraph(tf_imp_r, "• Dockerized Infrastructure:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=1)
    add_paragraph(tf_imp_r, "  Provided `docker-compose.yml` defining React served by Nginx alongside Flask + Gunicorn on production ports.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT)

    # -------------------------------------------------------------
    # SLIDE 9: Roadmap & Future Work (Light Theme)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    apply_bg(slide9, LIGHT_BG)
    add_header(slide9, "Roadmap & Future Plan")
    
    # Left Card: Short & Medium Term
    draw_card(slide9, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.7), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb_left = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.1), Inches(4.3))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    add_paragraph(tf_l, "Short & Medium Term (v3.5)", font_name="Georgia", font_size=18, color=ACCENT_TEAL, bold=True, space_after=16)
    
    add_paragraph(tf_l, "• Frontend Unit Testing:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_l, "  Set up Vitest and React Testing Library to validate UI routing, timer limits, and state bindings.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=12)
    
    add_paragraph(tf_l, "• Automated CI Pipeline:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_l, "  Configure GitHub Actions workflows to lint, auto-format, and execute backend/frontend tests on commits.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=12)
    
    add_paragraph(tf_l, "• Controlled User Study (N≈10):", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_l, "  Conduct trial runs with students to measure quantitative score improvements across consecutive practice sessions.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT)

    # Right Card: Long Term Vision
    draw_card(slide9, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.7), CARD_BG_LIGHT, CARD_BORDER_LIGHT)
    tb_right = slide9.shapes.add_textbox(Inches(7.23), Inches(2.0), Inches(5.1), Inches(4.3))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    add_paragraph(tf_r, "Long Term Vision (v4.0)", font_name="Georgia", font_size=18, color=ACCENT_TEAL, bold=True, space_after=16)
    
    add_paragraph(tf_r, "• User Authentication & Cloud profiles:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_r, "  Replace offline JSON file tracking with PostgreSQL databases and OAuth user profiles (Google / GitHub).", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=12)
    
    add_paragraph(tf_r, "• Collaborative Mock Sessions:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_r, "  Enable peer-to-peer live mock sessions alongside AI scoring to blend automated and human critiques.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT, space_after=12)
    
    add_paragraph(tf_r, "• Native Mobile Applications:", font_name="Calibri", font_size=15, color=TEXT_LIGHT_BG, bold=True, space_after=2)
    add_paragraph(tf_r, "  Compile responsive layouts using React Native for Android and iOS devices, enabling practice on-the-go.", font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT)

    # -------------------------------------------------------------
    # SLIDE 10: Viva-Voce / Q&A Prep (Light Theme)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    apply_bg(slide10, LIGHT_BG)
    add_header(slide10, "Common Viva-Voce Questions & Defense")
    
    questions = [
        {
            "q": "Q1: What is the core technical contribution of this work?",
            "a": "It combines Named Entity Recognition (NER via spaCy) with Large Language Models (Gemini/HF) to deliver contextual question-response-evaluation loops, paired with learning analytic charts, completely free of generic static templates."
        },
        {
            "q": "Q2: How does the system work when the AI APIs are offline or missing keys?",
            "a": "It implements a robust Fallback Mode. Questions load from a local bank of 100+ vetted files. Evaluations run through a regex-based heuristic model measuring length, structure, and keyword density."
        },
        {
            "q": "Q3: How are user sessions persisted, and what are the security measures?",
            "a": "Data is saved locally in data/sessions.json for the demo. Rate-limiting guards endpoints, strict CSP headers block code injections, and CORS rules prevent cross-site resource access."
        }
    ]
    
    qw = Inches(11.833)
    qh = Inches(1.3)
    y_coords = [Inches(1.8), Inches(3.3), Inches(4.8)]
    
    for idx, q_item in enumerate(questions):
        c_y = y_coords[idx]
        draw_card(slide10, Inches(0.75), c_y, qw, qh, CARD_BG_LIGHT, CARD_BORDER_LIGHT)
        
        tb = slide10.shapes.add_textbox(Inches(0.95), c_y + Inches(0.1), qw - Inches(0.4), qh - Inches(0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        add_paragraph(tf, q_item["q"], font_name="Georgia", font_size=14, color=ACCENT_TEAL, bold=True, space_after=2)
        add_paragraph(tf, q_item["a"], font_name="Calibri", font_size=13, color=TEXT_MUTED_LIGHT)

    # -------------------------------------------------------------
    # SLIDE 11: Conclusion & Thank You (Dark Theme)
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    apply_bg(slide11, DARK_BG)
    
    # Left accent colored strip
    strip11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    strip11.fill.solid()
    strip11.fill.fore_color.rgb = ACCENT_CYAN
    strip11.line.fill.background()
    
    title_box11 = slide11.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf11 = title_box11.text_frame
    tf11.word_wrap = True
    
    p_thank = add_paragraph(tf11, "Thank You!", font_name="Georgia", font_size=54, color=TEXT_DARK_BG, bold=True, space_after=12)
    p_thank.alignment = PP_ALIGN.CENTER
    
    p_disc = add_paragraph(tf11, "Questions & Discussion", font_name="Calibri", font_size=20, color=ACCENT_CYAN, space_after=24)
    p_disc.alignment = PP_ALIGN.CENTER
    
    line11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.666), Inches(4.3), Inches(2.0), Inches(0.04))
    line11.fill.solid()
    line11.fill.fore_color.rgb = ACCENT_TEAL
    line11.line.fill.background()
    
    info_box11 = slide11.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.5))
    tf_info11 = info_box11.text_frame
    tf_info11.word_wrap = True
    p_proj = add_paragraph(tf_info11, "AI-Based Resume Driven Interview System v2.0", font_name="Calibri", font_size=15, color=TEXT_DARK_BG, bold=True)
    p_proj.alignment = PP_ALIGN.CENTER
    p_sub_proj = add_paragraph(tf_info11, "Build. Practice. Excel.", font_name="Calibri", font_size=13, color=TEXT_MUTED_DARK)
    p_sub_proj.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_filename = "ai-interview-system-presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'")

if __name__ == "__main__":
    create_presentation()
